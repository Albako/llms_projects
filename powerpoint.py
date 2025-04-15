from pptx import Presentation

presentation = Presentation()
slide_layout = presentation.slide_layouts[1] # title and content layout
slide = presentation.slides.add_slide(slide_layout)
slide.placeholders[0].text = "My Title"
slide.placeholders[1].text = "Apple\nOranges"

presentation.save("sample.pptx")
